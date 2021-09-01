from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR_INDEX
from scripts.webscrapping.test_webscrap import get_blog_data
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData ,ChartData
from pptx.enum.chart import XL_CHART_TYPE,XL_TICK_MARK,XL_LEGEND_POSITION,XL_DATA_LABEL_POSITION

def get_cell_text(table,row,col):
  cell = table.cell(row,col)
  paragraphs = cell.text_frame.paragraphs
  for paragraph in paragraphs:
    for run in paragraph.runs:
      return run.text
      

def delete_slides(presentation, index):
  xml_slides = presentation.slides._sldIdLst  
  slides = list(xml_slides)
  xml_slides.remove(slides[index]) 

from pptx.util import Inches,Pt
import pandas as pd
from scripts.pandas_files.itter_rows import get_tuple_list
# title_id=[]
# slide=[]
# placeholders_id=[]
prs=Presentation("scripts/ppt_generation/ppts/input_4.pptx")
# for i in range(len(prs.slides)):
#   for shape in prs.slides[i].placeholders:
#     print('%d %d %s' % (i,shape.placeholder_format.idx, shape.name))
    
# for index,slide in enumerate(prs.slides):
# for shape in prs.slides[5].shapes.placeholders:
#   print('%d %s' % (shape.placeholder_format.idx, shape.name))

shape_name={}
funct_dict={}
for index,i in  enumerate(prs.slides[11].shapes):
  # print(f"id:{index} , name : {i.name}, shape_type :{i.shape_type} ")
  shape_name.update({i.name:[index,str(i.shape_type)]})
  funct_dict.update({i.name:index})
  
# print(prs.slides[0].shapes.title)
# print(shape_name)
# alt=[]
# for shape_name in prs.slides[0].shapes:
#       alt.append('%s' % (shape_name.name))

# for ind in alt:
#   shape_name=alt.index(ind)
#   print(shape_name)

# print(table.cell(1,1).text)

# print(funct_dict)
# data={"Addition":["145","190","133","94","179","206","78"],
# "Removals":["-1.24%","-11.3%","-4.5%","-1.21%","-1.05%","-9.0%","-5.1%"]}
# df=pd.DataFrame(data)
# print(df["Addition"].dtype)

# print(prs.slides[19].shapes[6])
  
# for shape in prs.slides[19].placeholders:
#   print('%d %s' % ( shape.placeholder_format.idx, shape.name))

  
  

df=pd.read_csv("scripts/ppt_generation/chart.csv")

print(df.to_dict("list"))
# df= pd.read_csv('scripts/pandas_files/csvfiles/uscities.csv')


# os.path.join(os.path.expanduser("~"),'config/sqlcred.cfg
def table_data():
  prs=Presentation("scripts/ppt_generation/themes/sharecare_temp.pptx")

  slide = prs.slides[19]
  shape = slide.shapes

  table = shape[7].table
  try:
    
    # print(table.cell(0,0).textframe.paragraphs.runs.text)
    # print(len(table.rows))
    # print(len(table.columns))
    table_data=dict()
    for j in range(len(table.columns)):
      table_data.update({get_cell_text(table,0,j):[]})
      
    # print(cols)
    # print(table.cell(0,0).text_frame.paragraphs.runs.text)
    # data={table.cell(0,0).text_frame.paragraphs.runs.text}
    for i in range(1,len(table.rows)):
      for j in range(0,len(table.columns)):
        table_data[get_cell_text(table,0,j)].append(get_cell_text(table,i,j))
    
    print(table_data)
  except Exception as error:
    print(error)
    