from scripts.webscrapping.test_webscrap import get_blog_data
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData 
from pptx.enum.chart import XL_CHART_TYPE,XL_TICK_MARK,XL_LEGEND_POSITION,XL_DATA_LABEL_POSITION

from pptx.util import Inches,Pt
import pandas as pd
from scripts.pandas_files.itter_rows import get_tuple_list

def get_table_data():
    path="scripts\pandas_files\csvfiles\students.csv"
    df=pd.read_csv(path,sep="|")
    return df
    
def main():
    prs = Presentation()

    ## SLIDE 1 
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "PPT generation using python"
    slide.placeholders[1].text = "By Abhinav"


    ## SLIDE 2
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left =Inches(1)
    top=Inches(1.3) 
    width=Inches(8)
    height = Inches(7)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "This is text inside a textbox"
    p = tf.add_paragraph()
    p.text = "This is a second paragraph that's bold"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "This is a third paragraph that's big"
    p.font.size = Pt(40)


    ## SLIDE 3
    third_slide=prs.slide_layouts[1]
    slide=prs.slides.add_slide(third_slide)
    shape=slide.shapes
    body_shape=shape.placeholders[1]
    shape.title.text="adding bullets"

    text_frame=body_shape.text_frame
    text_frame.text="topic"

    para=text_frame.add_paragraph()
    para.text="subtopic 1"
    para.level=1

    para=text_frame.add_paragraph()
    para.text="subtopic 2"
    para.level=1

    ## slide 4
    img_path="D:\Abhinav\download.jpeg"
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = Inches(1)
    pic = slide.shapes.add_picture(img_path, left, top)

    left = Inches(5)
    height = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)


    ## slide 5 
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'Adding a Table'
    left = top = Inches(2.0)
    width = Inches(9.0)
    height = Inches(0.8)
    df=get_table_data()
    list_tup=get_tuple_list(df)
    list_tup.insert(0,tuple(df.columns))
    cols=len(list_tup[0])
    row=len(list_tup)
    table = shapes.add_table(row, cols, left, top, width, height).table
    for i in range(row):
        for j in range(cols):
            table.cell(i, j).text=str(list_tup[i][j])
    
    ##SLIDE 6
    slide = prs.slides.add_slide(title_only_slide_layout)
    slide.shapes.title.text = "Graphical Representation"
    graph_data=CategoryChartData()
    graph_data.categories=["A","B","C"]
    graph_data.add_series('Q1 Sales', (19.2, 21.4, 16.7))
    graph_data.add_series('Q2 Sales', (22.3, 28.6, 15.2))
    graph_data.add_series('Q3 Sales', (20.4, 26.3, 14.2))
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5) 
  
    chart=slide.shapes.add_chart( XL_CHART_TYPE.COLUMN_CLUSTERED, x,y, cx, cy, graph_data ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.font.size = Pt(13)
    # data_labels.number_format = '0%'
    data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END
    
    ##SLIDE 7
    slide = prs.slides.add_slide(title_only_slide_layout)
    slide.shapes.title.text = "Pie chart Representation"
    chart_data = CategoryChartData()
    chart_data.categories = ['West', 'East', 'North', 'South', 'Other']
    chart_data.add_series('Series 1', (0.135, 0.324, 0.180, 0.235, 0.126))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END
    
    prs.save('scripts/ppt_generation/test.pptx')
    
if __name__=="__main__":
    main()