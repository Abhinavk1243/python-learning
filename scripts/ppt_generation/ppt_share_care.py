from scripts.webscrapping.test_webscrap import get_blog_data
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData ,ChartData
from pptx.enum.chart import XL_CHART_TYPE,XL_TICK_MARK,XL_LEGEND_POSITION,XL_DATA_LABEL_POSITION
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.util import Inches,Pt
import numpy as np
import pandas as pd
from pptx.enum.text import PP_ALIGN
from scripts.pandas_files.itter_rows import get_tuple_list
from pptx.table import _Cell
import copy
import os
from pptx.dml.color import RGBColor

def add_hyperlink_to_shape(shape,shape_id,url):
    shape[shape_id].click_action.hyperlink.address = url
    
def delete_slides(presentation, index):
    xml_slides = presentation.slides._sldIdLst  
    slides = list(xml_slides)
    xml_slides.remove(slides[index]) 

def replace_picture(slide, shape_no, img):
    
    """
    Replaces a picture as a placeholder on the template

    Arguments: -- slide index, shapename and image to be replaced with.
    Returns:-- None, as it simply replaces the image on the template
    """
    old_picture = slide.shapes[shape_no] 
    new_picture = slide.shapes.add_picture(img, top=Inches(slide.shapes[shape_no].top.inches), 
                                        left=Inches(slide.shapes[shape_no].left.inches),height=Inches(slide.shapes[shape_no].height.inches), width=Inches(slide.shapes[shape_no].width.inches))
    old_pic = old_picture._element
    new_pic = new_picture._element
    old_pic.addnext(new_pic) 
    old_pic.getparent().remove(old_pic)

def remove_row(table,del_rows,len_rows,row_count):
    tbl = table._tbl
    a={}
    for k in range(len_rows,row_count):
        key = k
        value = del_rows[k]
        a[key] = value
        tr = a[key]._tr
        tbl.remove(tr)
        
def add_row(table,row_count,len_rows):
    copy_idx=0
    lst=[]
    for k in range(row_count,len_rows):
        k = copy.deepcopy(table._tbl.tr_lst[copy_idx])
        lst.append(k)
    for rows in lst:
        for tc in rows.tc_lst:
            cell = _Cell(tc, rows.tc_lst)
        table._tbl.append(rows)
        return (table)
def replace_table(shape,shape_id,data):
    table=shape[shape_id].table
    row_count = len(table.rows)
    df=pd.DataFrame(data)
    col=list(df.columns)
    if len(df)<row_count:
        for i in range(0,len(df)):
            for j in range(0,len(col)):
                if (df[col[j]]).dtype == object :
                    shape[shape_id].table.cell(i+1,j).text=str(df[col[j]][i])
                    paragraph = shape[shape_id].table.cell(i+1,j).text_frame.paragraphs[0]
                    paragraph.font.size = Pt(9)
                else:
                    value=int(df[col[j]][i])
                    shape[shape_id].table.cell(i+1,j).text=str(value)
                    paragraph = shape[shape_id].table.cell(i+1,j).text_frame.paragraphs[0]
                    paragraph.font.size = Pt(9)
        del_rows=[]
        len_rows=len(df)+1
        for i in range(0,row_count):
            del_rows.append(table.rows[i])
        remove_row(table,del_rows,len_rows,row_count)  
    else:
        for i in range(0,row_count-1):
            for j in range(0,len(col)):
                if (df[col[j]]).dtype == object :
                    shape[shape_id].table.cell(i+1,j).text=str(df[col[j]][i])
                    paragraph = shape[shape_id].table.cell(i+1,j).text_frame.paragraphs[0]
                    paragraph.font.size = Pt(9)
                else:
                    value=int(df[col[j]][i])
                    shape[shape_id].table.cell(i+1,j).text=str(value)
                    paragraph = shape[shape_id].table.cell(i+1,j).text_frame.paragraphs[0]
                    paragraph.font.size = Pt(9)
        len_rows=len(df)
        table=add_row(table,row_count,len_rows)
        for i in range(row_count,len_rows-1):
            for j in range(0,len(col)):
                table.cell(i,j).text=str(df[col[j]][i])
                paragraph = table.cell(i,j).text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)

def replace_chart(slide,shape_id,data):
    """[summary]

    Args:
        slide ([type]): [description]
        shape_id ([type]): [description]
        df (Dataframe): [description]
    """
    df=pd.DataFrame(data)
    df=df.replace(np.nan, 0)
    chart = slide.shapes[shape_id].chart
    chart_data = ChartData()
    category=list(df.columns)
    del category[0]
    chart_data.categories =list(df["category"])
    series=list(df.columns)
    del series[0]
    series_count=len(chart.series)
    if series_count==len(df.columns)-1:  
        for index,i in enumerate(series):
            chart_series=chart.series[index]
            df[i]=df[i].replace('%', '', regex=True)
            if df[i].dtype==object:
                chart_data.add_series(chart_series.name,tuple(df[i]), number_format='0.0%')
            else:
                chart_data.add_series(chart_series.name,tuple(df[i]),number_format='#,##0')
        chart.replace_data(chart_data)
    

def replace_text(shape,shape_name,shape_id,shape_value):
    for i in range(len(shape_name)):
        txbox=shape[shape_id[i]]
        textframe = txbox.text_frame
        p = textframe.paragraphs[0]
        run = p.add_run()
        run.text = shape_value[i]
        font = run.font
        font.name = 'Consolas'
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        font.size = Pt(28)
        
def replace_funnel_text(shape,shape_name,shape_id,shape_value):
    for i in range(len(shape_name)):
        txbox=shape[shape_id[i]]
        textframe = txbox.text_frame
        p = textframe.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = shape_value[i]
        font = run.font
        font.color.rgb = RGBColor(255, 255, 255)
        # font.name = 'Calibri'
        # font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
        font.size = Pt(10)
        


def main():
    prs=Presentation("scripts/ppt_generation/ppts/input_4.pptx")

    ## program overview
    slide=prs.slides[3]
    shape=slide.shapes
    shape_name={'regn_percent': 21, 'reslage_test_completion_percent': 22, 'incentive_participation': 23, 'monthly_user': 24,
                '90day_engagement': 25, 'challenge_ptpn': 26, 'content_clicks': 27, 'lifestyle_mgmt_enrollment': 28, 
                '1+condition_risk': 29, '1+lifestyle_risk': 30, 'net_risk_change': 31, 'preventive_risk_change': 32, 
                'lifestyle_risk_change': 33, 'biometric_risk_change': 34}
    shape_id=list(shape_name.values())
    shape_value=["21.5%","4.5%","30%","45%","27%","29.4%","67.9","20%","45.3%","97.4%","-12.3%","-11.8%","-12.1","-5.1%"]


    replace_text(shape,shape_name,shape_id,shape_value)


    ## eligibility trends 

    slide=prs.slides[6]
    shape=slide.shapes
    shape_nameidx={'eligible_members': 7 ,'current_q_vs_prior_q':8}
    shape_value=["4678","-5.6%"]
    shape_id=list(shape_nameidx.values())
    replace_text(shape,shape_nameidx,shape_id,shape_value)

    data={"category":["2021-02","2020-3","2020-4","2020-5","2020-6","2020-3","2020-3"],
        "overall":[4300,4390,4270,4314,4479,4596,4678]}

    shape_id=3
    replace_chart(slide,shape_id,data)

    data={"category":["2021-02","2021-03","2021-04","2021-05","2021-06","2021-07","2021-08"],
    "Addition":[145,190,133,94,179,206,78],
    "Removals":[-124,-113,-45,-121,-105,-90,-51]}

    shape_id=4

    replace_chart(slide,shape_id,data)


    ##Registration and RealAge Test Completion 
    slide=prs.slides[8]
    shape=slide.shapes
    shape_name={'realage_eligible': 15, 'realage_regd_2': 16,
                'regd_members': 17, 'distinct_test_completers': 18, 'distinct_test_completers_curr_year': 19, 'regn_rate': 20,
                'realage_regd': 21, 'current_q_vs_prior_q_regn': 22}
    shape_value=["6.2%","19.7%","1123","892","231","23.4%","4.1%","79%"]
    shape_id=list(shape_name.values())
    replace_text(shape,shape_name,shape_id,shape_value)



    data={"category":["2021-02","2020-3","2020-4","2020-5","2020-6","2020-3","2020-3"],
        "Eligible_member":[4300,4390,4270,4314,4479,4596,4678],
        "registeration_rate":[".187%",".196%",".219%",".202%",".224%",".231%",".213%"],
        "RAT_comp_rate":["0.047%","0.036%","0.019%","0.022%","0.044%","0.031%","0.043%"]
        }

    shape_id=13

    replace_chart(slide,shape_id,data)

    # Registration and RealAge Completion by Group 2,3,7 chart

    slide=prs.slides[9]
    shape=slide.shapes

    data={"category":["employes","spouse","others"],
                    "eligible":[4512,771.5,5.56],
                    "reg_bef_period":[".265%",".0556%","0.0%"],
                    "reg_after_period":["0.21%","0.04%",".0%"]}
    shape_id=7
    replace_chart(slide,shape_id,data)
    data={'category': ['Female', 'Male'], 'Eligible': [2004, 2514], 'Registration_Rate': ["0.3273","0.2033"], 'RAT_Completion Rate': ["0.2665","0.161"]}
    replace_chart(slide,2,data)


    shape_id=3
    data={'category': ['0-17', '18-24', '25-34', '35-44', '45-54', '55-64', '65+'],
    'Eligible': [0, 361, 1228, 1072, 1032, 720, 105], 
    'Registration_rate': ['0.0%', '4.43', '23.13', '29.29', '31.3', '28.33', '24.76'],
    'RAT Completion Rate': ['0.0%', '2.77', '18.24', '24.44', '26.07', '22.64', '20']}

    replace_chart(slide,shape_id,data)

    ##Risk Analysis Summary1

    slide=prs.slides[12]
    shape=slide.shapes
    table=shape[2].table

    data={
        "lifestyle_risk":['Appointment Adherence','Binge Drinking','Depression (PHQ-2)','Diet','Excessive Drinking','Medication Adherence',
                          'Overweight','Physical Activity','Preventive Care','Sedentary Lifestyle','Sleep','Stress','Tabacoo','Weight Taining'],
        "book_of_bussiness":["1.5%","3.0%","5.6%","62.0%","3.4%","5.4%","33.3%","73.5%","60.9%","58.3%","7.4%","10.4%","5.4%","47.0%"],
        "client":["2.4%","5.2%","6.9%","68.0%","2.5%","4.7%","35.2%","68.5%","61.5%","60.7%","8.1%","8.5%","4.6%","43.6%"]
        }
    replace_table(shape,2,data)
    
    ## Biometric Screening Participation Slide
    slide=prs.slides[13]
    shape=slide.shapes
    shape_name={'eligible_members_biometric': 10,'screening_participants': 11,'participation_rate_biometric': 12}
    shape_value=["4324","341","6.5%"]
    shape_id=list(shape_name.values())
    replace_text(shape,shape_name,shape_id,shape_value)
    data={"category":["male","female"],
          "participents":[63.4,36.6]}
    shape_id=9
    replace_chart(slide,shape_id,data)
    data={"category":["January","feburary","March","April","May","June","July","August","September","October","Novemeber","December"],
     'prior_year': [10, 46, 52, 56, 61, 63, 63, 63, 63, 69, 73, 88],
     'current_year': [12.0, 19.0, 26.0, 32.0, 41.0, 57.0, None,None, None,None,None, None]}
    
    shape_id=8
    replace_chart(slide,shape_id,data)
    
    ##Demographics and Data Quality
    slide=prs.slides[7]
    shape=slide.shapes
    img=os.path.join(os.path.expanduser("~"),'flask-learning/map_img/map.jpeg')
    replace_picture(slide,12,img)
    data={
          'category': ['0-17', '18-24', '25-34', '35-44', '45-54', '55-64', '65+'],
          'eligible_members': ['0.0%', '8%', '27.2%', '23.7%', '22.8%', '15.9%', '2.3%']
        }
    shape_id=3
    replace_chart(slide,shape_id,data)
    shape_name={'avg_age': 10,'missing_email': 11,'missing_phone': 13}
    shape_id=list(shape_name.values())
    shape_value=["41.3","74.1%","0.1%"]
    replace_text(shape,shape_name,shape_id,shape_value)
    data={"category":["Female","Male"],
          "Members" :["44.36","55.64"]}
    shape_id=7
    replace_chart(slide,shape_id,data)
    # Biometric/Clinical Screening Results
    slide=prs.slides[14]
    shape=slide.shapes
    data={'category': ['Total Cholesterol/HDL Ratio', 'Glucose (Fasting)', 'Glucose', 'Cholesterol - Total', 'Cholesterol - LDL (Bad)', 'Cholesterol - HDL (Good)', 'BMI (Body Mass Index)', 'Blood Pressure - Systolic (Upper)', 'Blood Pressure - Diastolic (Lower)'],
          'Series_1%': ['0.566%', '0.715%', '0.509%', '0.667%', '0.392%', '0.399%', '0.293%', '0.467%', '0.585%'],
          'Series_2%': ['0.0%', '0.253%', '0.0%', '0.279%', '0.557%', '0.217%', '0.351%', '0.39%', '0.288%'],
          'Series_3': ['0.434%', '0.032%', '0.491%', '0.054%', '0.052%', '0.384%', '0.356%', '0.143%', '0.127%']}
    replace_chart(slide,9,data)
    
    ## funnel chart
    
    slide=prs.slides[24]
    shape=slide.shapes
    shape_name={'Targeted1': 14, 'Callable1': 15,  'Attempted1': 16, 'Enrolled1': 17}
    shape_value=["798","768","575","124"]
    shape_id=list(shape_name.values())
    replace_funnel_text(shape,shape_name,shape_id,shape_value)
    
    
    slide=prs.slides[37]
    shape=slide.shapes
    data={'Risk Name': ['BMI (High)', 'Blood Pressure (High)', 'PHQ2 - Depressed/Little Interest'],
          'Number of Risk Points at T1': ['114', '35', '24']}
    replace_table(shape,6,data)
    
    
    ## real age slide 
    slide=prs.slides[11]
    shape=slide.shapes
    
    shape_name={'avg_realage_delta': 7, 'members_with_realage_delta': 8, 'distinct_realage_test_completers': 9}
    shape_id=list(shape_name.values())
    shape_value=["-2.7","13%","875"]
    replace_text(shape,shape_name,shape_id,shape_value)
    add_hyperlink_to_shape(shape,7,"https://github.com/Abhinavk1243/python-learning")
    
    data={'category': ['LT -5', '-4.9 to -3', '-2.9 to -1', '-.9 to 1', '1.1 to 2.9', '3 to 5', 'GT 5'],
     'Series_1': [20.9, 30.3, 24.6, 12.1, 6.2, 2.9, 2.9]}
    replace_chart(slide,10,data)
    
    ## challenge participation sponsor initiated
    slide=prs.slides[20]
    shape=slide.shapes
    shape.title.text="Challenge Participation – Sponsor Initiated"
    
    shape_name={'challeng_ptpn_rate': 12, 'challenge_joins': 13, 'challenge_completions': 14,'challenge_completion_rate': 15}
    shape_id=list(shape_name.values())
    shape_value=["33.8%","2377","1385","54%"]
    
    replace_text(shape,shape_name,shape_id,shape_value)
        
    data={'Challenge Type': ['Individual'], 'Challenge Participations': ['2,377'], 'Challenge Completes': ['1,385'], 'Challenge Completion Rate': ['58.3%']}
    replace_table(shape,7,data)
    
    data={'category': ['2020-04', '2020-052', '2020-06', '2020-07', '2020-08', '2020-09', '2020-10', '2020-11', '2020-12', '2021-01', '2021-02', '2021-03', '2021-04', '2021-05', '2021-06'], 
          'Challenge Participations': [185, 71, 164, 191, 0, 0, 0, 0, 0, 161, 166, 168, 149, 117, 90], 
          'Challenge Completes': [134, 178, 113, 148, 0, 0, 0, 0, 0, 78, 104, 87, 70, 56, 49]}
    
    replace_chart(slide,6,data)
    
    
    
    delete_slides(prs,4)
    delete_slides(prs,11)
    prs.save('scripts/ppt_generation/ppts/output_4.pptx')
    
    
if __name__=="__main__":
    main()