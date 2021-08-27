from scripts.webscrapping.test_webscrap import get_blog_data
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData ,ChartData
from pptx.enum.chart import XL_CHART_TYPE,XL_TICK_MARK,XL_LEGEND_POSITION,XL_DATA_LABEL_POSITION

from pptx.util import Inches,Pt
import pandas as pd
from scripts.pandas_files.itter_rows import get_tuple_list



prs=Presentation("scripts/ppt_generation/two_axix_graph.pptx")
# chart = prs.slides[0].shapes[0].chart
# chart_data = ChartData()
# chart_data.categories =["A","B","C"]
# chart_data.add_series("sales1",(2.34,3.3,2))
# chart_data.add_series("sales2",(5,6.3,1))
# chart_data.add_series("sales3",(4.34,6.3,5))
# chart.replace_data(chart_data)

# prs.save("test_2axis.pptx")
# print(shape)
# slide=prs.slides[0].shapes
# print(slide)
# for slide in prs.slides:
#     for shape in slide.placeholders:
#         print('%d %d %s' % (prs.slides.index(slide), shape.placeholder_format.idx, shape.name))


for i in prs.slides[1].shapes:
    print(i.name)
    