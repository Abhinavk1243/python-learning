from scripts.webscrapping.test_webscrap import get_blog_data
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData ,ChartData
from pptx.enum.chart import XL_CHART_TYPE,XL_TICK_MARK,XL_LEGEND_POSITION,XL_DATA_LABEL_POSITION

from pptx.util import Inches,Pt
import pandas as pd
from scripts.pandas_files.itter_rows import get_tuple_list

prs=Presentation("scripts/ppt_generation/themes/column2_line_temp.pptx")
prs.slides[0].shapes.title.text = "2 axis by rendering data"
# for shape in prs.slides[0].placeholders:
#        print('%d %s' % (shape.placeholder_format.idx, shape.name))
# title=prs.slides[0].shapes[0].placeholders[0]
# title.text="2 axis"
chart = prs.slides[0].shapes[1].chart
chart_data = ChartData()
chart_data.categories =["A","B","C","D","E"]
chart_data.add_series("sales1",(234,453,100,600,540))
chart_data.add_series("sales2",(560,290,680,220,320))
chart_data.add_series("sales1",(834,630,300,430,100))
chart.replace_data(chart_data)
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.LEFT
chart.legend.include_in_layout = False
    
chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.font.size = Pt(13)
# data_labels.number_format = '0%'
data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END

prs.save("test_axis.pptx")